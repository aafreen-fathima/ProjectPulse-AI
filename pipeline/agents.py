"""
LangGraph orchestration — the spine of the weekly report pipeline.

Graph:

    collect ─▶ score ─▶ analyse_gaps ─▶ visualise ─▶ narrate ─▶ format ─▶ distribute

Each node receives and mutates a `PipelineState` dict.

Run locally:
    python -m pipeline.agents --org-id org_demo
"""

from __future__ import annotations

import argparse
import io
import logging
import os
from datetime import date, timedelta
from typing import TypedDict

import anthropic
import pandas as pd
from langgraph.graph import END, StateGraph
from pptx import Presentation
from pptx.util import Inches
from resend import Resend
from supabase import create_client

from .gap_analysis import compute as compute_variance
from .risk_model import FEATURES, score_projects, synthetic_training_set, train
from .visualiser import budget_waterfall, risk_heatmap, risk_trend, buf

log = logging.getLogger("agents")

ANTHROPIC_MODEL = os.getenv("ANTHROPIC_MODEL", "claude-sonnet-4-6")


class PipelineState(TypedDict, total=False):
    org_id: str
    period_start: str
    period_end: str
    projects: pd.DataFrame
    milestones: pd.DataFrame
    risk_scores: list
    variance: list
    charts: dict[str, bytes]
    narrative: str
    pptx_bytes: bytes
    distributed: bool


# ────────────────────────────────────────────────────────────
# Nodes
# ────────────────────────────────────────────────────────────
def collect(state: PipelineState) -> PipelineState:
    sb = create_client(os.environ["NEXT_PUBLIC_SUPABASE_URL"], os.environ["SUPABASE_SERVICE_ROLE_KEY"])
    projects = pd.DataFrame(sb.table("projects").select("*").eq("org_id", state["org_id"]).execute().data)
    milestones = pd.DataFrame(sb.table("milestones").select("*").eq("org_id", state["org_id"]).execute().data)
    log.info("collected %d projects / %d milestones", len(projects), len(milestones))
    return {**state, "projects": projects, "milestones": milestones}


def score(state: PipelineState) -> PipelineState:
    # Build feature frame from projects + milestones. In production these
    # join on Jira and Slack signal too.
    proj = state["projects"].copy()
    feat = pd.DataFrame(
        {
            "project_id":        proj["id"],
            "task_age":          (pd.Timestamp.utcnow() - pd.to_datetime(proj["start_date"])).dt.days / 7,
            "velocity_trend":    proj["sentiment"].fillna(0) * 0.5,        # placeholder coupling
            "dependency_depth":  2,
            "sentiment_score":   proj["sentiment"].fillna(0),
            "budget_burn_rate":  (proj["spent_usd"] / proj["budget_usd"].replace(0, pd.NA)).fillna(0),
            "days_to_milestone": 14,
            "schedule_variance": 0,
        }
    )

    # Bootstrap from synthetic data until real labelled history exists.
    model = train(synthetic_training_set(), log_to_mlflow=False)
    scores = score_projects(model, feat)
    log.info("scored %d projects", len(scores))
    return {**state, "risk_scores": scores}


def analyse_gaps(state: PipelineState) -> PipelineState:
    variance = compute_variance(state["projects"], state["milestones"], baseline_count={})
    return {**state, "variance": variance}


def visualise(state: PipelineState) -> PipelineState:
    # Build a fake weekly history series (in prod: pulled from `risk_history`).
    weeks = pd.date_range(end=pd.Timestamp.utcnow(), periods=8, freq="W").strftime("%Y-%m-%d").tolist()
    rows = []
    for w in weeks:
        for s in state["risk_scores"]:
            rows.append({"week": w, "project_code": s.project_id[:8], "score": s.score})
    scores_df = pd.DataFrame(rows)

    history = scores_df.groupby("week", as_index=False)["score"].mean().rename(columns={"score": "avg_score"})

    charts = {
        "heatmap":   risk_heatmap(scores_df),
        "trend":     risk_trend(history),
        "waterfall": budget_waterfall([{"project_code": v.project_code, "budget_pct": v.budget_pct} for v in state["variance"]]),
    }
    return {**state, "charts": charts}


def narrate(state: PipelineState) -> PipelineState:
    """Claude drafts the 600-word executive narrative."""
    client = anthropic.Anthropic(api_key=os.environ["ANTHROPIC_API_KEY"])

    # Compress signal into a tight prompt
    proj_lines = "\n".join(
        f"- {row.code} ({row['name']}): status={row.status}, rag={row.rag}, spent={row.spent_usd}/{row.budget_usd}"
        for _, row in state["projects"].iterrows()
    )
    risk_lines = "\n".join(
        f"- {s.project_id[:8]}: score={s.score:.2f}, sev={s.severity}, top={[f['feature'] for f in s.shap_features]}"
        for s in state["risk_scores"]
    )

    msg = client.messages.create(
        model=ANTHROPIC_MODEL,
        max_tokens=1200,
        system=(
            "You are the chief of staff to a PMO lead. Write a calm, decision-oriented "
            "executive summary (~600 words) for the week. Lead with the one question the "
            "leadership team needs to answer this week. Then: portfolio health, top 3 risks "
            "with mitigation, budget posture, recommended actions. No marketing fluff. "
            "No emojis. No bullet points longer than two lines."
        ),
        messages=[
            {
                "role": "user",
                "content": f"Period: {state['period_start']} to {state['period_end']}\n\n"
                           f"Projects:\n{proj_lines}\n\nRisks:\n{risk_lines}",
            }
        ],
    )
    narrative = "\n".join(b.text for b in msg.content if b.type == "text")
    return {**state, "narrative": narrative}


def format_deck(state: PipelineState) -> PipelineState:
    prs = Presentation()
    blank = prs.slide_layouts[6]

    # Title slide
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "ProjectPulse — Weekly Brief"
    s.placeholders[1].text = f"{state['period_start']} → {state['period_end']}"

    # Narrative slide
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Executive summary"
    tx = s.shapes.placeholders[1] if len(s.shapes.placeholders) > 1 else None
    if tx is None:
        # add a textbox manually
        tx = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.text = state["narrative"][:3000]

    # Chart slides
    for label, png in state["charts"].items():
        s = prs.slides.add_slide(blank)
        s.shapes.title_placeholder = None
        s.shapes.add_picture(buf(png), Inches(0.5), Inches(0.5), width=Inches(9))

    out = io.BytesIO()
    prs.save(out)
    return {**state, "pptx_bytes": out.getvalue()}


def distribute(state: PipelineState) -> PipelineState:
    sb = create_client(os.environ["NEXT_PUBLIC_SUPABASE_URL"], os.environ["SUPABASE_SERVICE_ROLE_KEY"])

    # Persist to Supabase Storage (bucket: reports)
    path = f"{state['org_id']}/{state['period_start']}.pptx"
    sb.storage.from_("reports").upload(path, state["pptx_bytes"], {"content-type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"})
    pptx_url = sb.storage.from_("reports").get_public_url(path)

    sb.table("reports").insert(
        {
            "org_id": state["org_id"],
            "period_start": state["period_start"],
            "period_end": state["period_end"],
            "narrative": state["narrative"],
            "pptx_url": pptx_url,
        }
    ).execute()

    # Email stakeholders
    resend = Resend(api_key=os.environ["RESEND_API_KEY"])
    recipients = [u["email"] for u in sb.table("users").select("email").eq("org_id", state["org_id"]).execute().data]
    if recipients:
        resend.emails.send(
            {
                "from": os.environ.get("RESEND_FROM_EMAIL", "reports@projectpulseai.com"),
                "to": recipients,
                "subject": f"ProjectPulse weekly brief — {state['period_start']}",
                "html": f"<p>This week's portfolio brief is ready.</p>"
                        f"<p><a href='{pptx_url}'>Download deck</a></p>"
                        f"<pre style='white-space:pre-wrap'>{state['narrative'][:1200]}…</pre>",
            }
        )
    return {**state, "distributed": True}


# ────────────────────────────────────────────────────────────
# Graph wiring
# ────────────────────────────────────────────────────────────
def build_graph():
    g = StateGraph(PipelineState)
    g.add_node("collect", collect)
    g.add_node("score", score)
    g.add_node("analyse_gaps", analyse_gaps)
    g.add_node("visualise", visualise)
    g.add_node("narrate", narrate)
    g.add_node("format", format_deck)
    g.add_node("distribute", distribute)

    g.set_entry_point("collect")
    g.add_edge("collect", "score")
    g.add_edge("score", "analyse_gaps")
    g.add_edge("analyse_gaps", "visualise")
    g.add_edge("visualise", "narrate")
    g.add_edge("narrate", "format")
    g.add_edge("format", "distribute")
    g.add_edge("distribute", END)
    return g.compile()


def run(org_id: str, period_start: str | None = None, period_end: str | None = None) -> dict:
    today = date.today()
    period_start = period_start or (today - timedelta(days=today.weekday())).isoformat()
    period_end = period_end or (date.fromisoformat(period_start) + timedelta(days=6)).isoformat()
    graph = build_graph()
    return graph.invoke({"org_id": org_id, "period_start": period_start, "period_end": period_end})


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(name)s %(message)s")
    parser = argparse.ArgumentParser()
    parser.add_argument("--org-id", required=True)
    args = parser.parse_args()
    final = run(args.org_id)
    print("done. distributed:", final.get("distributed"))
